import pdfplumber
from pptx import Presentation
from docx import Document
import pandas as pd
import os
from typing import List, Dict


def parse_file(file_path: str) -> List[Dict]:
    """
    Main dispatcher function to parse a supported file based on its extension.

    Args:
        file_path (str): Full path to the input file.

    Returns:
        List[Dict]: Parsed text or table data with metadata entries.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        entries = parse_pdf(file_path)
    elif ext == ".csv":
        entries = parse_csv(file_path)
    elif ext == ".pptx":
        entries = parse_pptx(file_path)
    elif ext == ".docx":
        entries = parse_docx(file_path)
    elif ext in [".txt", ".md"]:
        entries = parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

    return entries

def parse_pdf(file_path: str) -> List[Dict]:
    """
    Extracts non-table text and table rows from a PDF using pdfplumber.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        List[Dict]: Parsed entries with type: sentence or table_row.
    """
    results = []
    base_name = os.path.basename(file_path)

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            # Get table bounding boxes
            table_bboxes = [table.bbox for table in page.find_tables()]

            def not_within_bboxes(obj):
                """Exclude words inside any table bbox."""
                def in_bbox(bbox):
                    v_mid = (obj["top"] + obj["bottom"]) / 2
                    h_mid = (obj["x0"] + obj["x1"]) / 2
                    x0, top, x1, bottom = bbox
                    return x0 <= h_mid <= x1 and top <= v_mid <= bottom
                return not any(in_bbox(b) for b in table_bboxes)

            # Extract clean non-table text (punctuation preserved)
            text = page.filter(not_within_bboxes).extract_text()
            if text:
                results.append({
                    "text": text.strip(),
                    "source_file": base_name,
                    "page": page_num,
                    "type": "pdf_text"
                })

            # Extract tables separately
            for table_num, table in enumerate(page.extract_tables()):
                df = pd.DataFrame(table[1:], columns=table[0])
                for row_index, row_data in enumerate(table_to_kv_chunks(df)):
                    results.append({
                        "text": row_data,
                        "source_file": base_name,
                        "page": page_num,
                        "table": table_num + 1,
                        "row": row_index + 1,
                        "type": "table_row"
                    })

    return results

def parse_csv(file_path: str) -> List[Dict]:
    """
    Parses rows from a CSV file.

    Args:
        file_path (str): Path to the CSV file.

    Returns:
        List[Dict]: Parsed CSV rows as dictionaries.
    """
    df = pd.read_csv(file_path)
    base_name = os.path.basename(file_path)
    results = [
        {
            "text": row_data,
            "source_file": base_name,
            "row": i + 1,
            "type": "csv_row"
        }
        for i, row_data in enumerate(table_to_kv_chunks(df))
    ]
    return results

def parse_pptx(file_path: str) -> List[Dict]:
    """
    Parses text from PowerPoint slides.

    Args:
        file_path (str): Path to the .pptx file.

    Returns:
        List[Dict]: Slide-wise extracted text or table entries.
    """
    prs = Presentation(file_path)
    base_name = os.path.basename(file_path)
    results = []

    for i, slide in enumerate(prs.slides):
        slide_text_parts = []
        table_num = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and not shape.has_table:
                text = shape.text.strip()
                if text:
                    slide_text_parts.append(text)

            elif shape.has_table:
                table_num += 1
                try:
                    # Convert table to DataFrame
                    table = shape.table
                    data = []
                    for row in table.rows:
                        data.append([cell.text.strip() for cell in row.cells])
                    df = pd.DataFrame(data[1:], columns=data[0])  # First row = headers

                    # Chunk table into key-value sentences
                    for row_index, row_data in enumerate(table_to_kv_chunks(df)):
                        results.append({
                            "text": row_data,
                            "source_file": base_name,
                            "slide": i + 1,
                            "table": table_num,
                            "row": row_index + 1,
                            "type": "table_row"
                        })
                except Exception as e:
                    print(f"Error processing table on slide {i+1}: {e}")

        if slide_text_parts:
            results.append({
                "text": " ".join(slide_text_parts),
                "source_file": base_name,
                "slide": i + 1,
                "type": "pptx_slide"
            })

    return results

def parse_docx(file_path: str) -> List[Dict]:
    """
    Parses text from paragraphs in a .docx Word file.

    Args:
        file_path (str): Path to the .docx file.

    Returns:
        List[Dict]: Paragraph-wise extracted text.
    """
    results = []
    base_name = os.path.basename(file_path)
    doc = Document(file_path)

    # Extract each paragraph as a separate block
    for i, para in enumerate(doc.paragraphs, start=1):
        text = para.text.strip()
        if text:
            results.append({
                "text": text,
                "source_file": base_name,
                "paragraph": i,
                "type": "docx_paragraph"
            })

    # Extract tables into KV format
    for table_num, table in enumerate(doc.tables, start=1):
        data = []
        for row in table.rows:
            data.append([cell.text.strip() for cell in row.cells])
        if len(data) >= 2:
            df = pd.DataFrame(data[1:], columns=data[0])
        else:
            df = pd.DataFrame(data)

        for row_index, row_data in enumerate(table_to_kv_chunks(df)):
            results.append({
                "text": row_data,
                "source_file": base_name,
                "table": table_num + 1,
                "row": row_index + 1,
                "type": "table_row"
            })

    return results


def parse_txt(file_path: str) -> List[Dict]:
    """
    Parses text from paragraphs in a .txt file.

    Args:
        file_path (str): Path to the .txt file.

    Returns:
        List[Dict]: Paragraph-wise extracted text.
    """
    results = []
    base_name = os.path.basename(file_path)

    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read().strip()

    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]

    for i, para in enumerate(paragraphs, start=1):
        results.append({
            "text": para,
            "source_file": base_name,
            "paragraph": i,
            "type": "txt_paragraph"
        })

    return results

def table_to_kv_chunks(df: pd.DataFrame) -> List[str]:
    """
    Converts a DataFrame into a list of key-value formatted string chunks.

    For each row in the DataFrame:
    - Column names are treated as keys.
    - Cell values are treated as values.
    - Output format: "Column1: Value1; Column2: Value2; ..."
    - Columns with "Unnamed:" are treated as value-only (no key prefix).

    Args:
        df (pd.DataFrame): Input DataFrame to convert.

    Returns:
        List[str]: List of strings where each string represents one row in key-value format.
    """
    row_data = []
    for _, row in df.iterrows():
        kv_pairs = []
        for col in df.columns:
            col_name = str(col).strip() if col is not None else ""
            val = row[col]
            if isinstance(val, pd.Series):
                val = val.iloc[0]  # or apply logic to join values if needed

            cell_value = str(val).strip() if pd.notna(val) else ""

            if not col_name.startswith("Unnamed:"):
                kv_pairs.append(f"{col_name}: {cell_value}")
            else:
                kv_pairs.append(f"{cell_value}")
        row_data.append("; ".join(kv_pairs))
    return row_data